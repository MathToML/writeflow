from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # deep navy
BG_MEDIUM = RGBColor(0x1E, 0x29, 0x3B)      # slate 800
ACCENT = RGBColor(0x38, 0xBD, 0xF8)         # sky blue
ACCENT2 = RGBColor(0x34, 0xD3, 0x99)        # emerald
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)          # slate 300
DIM = RGBColor(0x94, 0xA3, 0xB8)            # slate 400
WARM = RGBColor(0xFB, 0xBF, 0x24)           # amber


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Calibri",
             space_before=0, line_spacing=1.3):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.line_spacing = Pt(font_size * line_spacing)
    return p


def add_shape_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent bar
add_shape_rect(slide, 0, 0, 0.15, 7.5, ACCENT)

add_text_box(slide, 1.5, 1.5, 10, 1.2, "CAMP", font_size=72,
             color=ACCENT, bold=True, font_name="Calibri")
tf = add_text_box(slide, 1.5, 3.0, 10, 0.8, "Calm Productivity",
                  font_size=36, color=WHITE, bold=False)
add_para(tf, "", font_size=12)
add_para(tf, "Don't manage your tasks. Let your tasks manage themselves.",
         font_size=20, color=DIM)

add_text_box(slide, 1.5, 5.8, 10, 0.5, "OTTD \u2014 One Thing To Do",
             font_size=16, color=DIM)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.6, 6, 0.8, "The Problem", font_size=40,
             color=ACCENT, bold=True)

tf = add_text_box(slide, 1.0, 1.7, 5.5, 4.5,
                  "We're drowning in productivity tools.",
                  font_size=24, color=WHITE, bold=True)
add_para(tf, "", font_size=10)
add_para(tf, "\u2022  20 tasks on the to-do list. Which one now?",
         font_size=20, color=LIGHT, space_before=8)
add_para(tf, "\u2022  Calendar, reminders, notes \u2014 scattered everywhere",
         font_size=20, color=LIGHT, space_before=8)
add_para(tf, "\u2022  More tools = more overhead = more stress",
         font_size=20, color=LIGHT, space_before=8)

# Quote box on right
add_shape_rect(slide, 7.2, 1.7, 5.3, 3.8, BG_MEDIUM)
tf = add_text_box(slide, 7.6, 2.0, 4.5, 3.2,
                  '\u201c I downloaded 3 apps to be productive.',
                  font_size=20, color=LIGHT)
add_para(tf, "  Now I spend 30 minutes a day", font_size=20, color=LIGHT)
add_para(tf, "  just organizing them. \u201d", font_size=20, color=LIGHT)
add_para(tf, "", font_size=14)
add_para(tf, "The irony:", font_size=22, color=WARM, bold=True,
         space_before=12)
add_para(tf, "Productivity tools are making us", font_size=22, color=WARM)
add_para(tf, "less productive.", font_size=22, color=WARM, bold=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: Just Dump It
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.6, 6, 0.8, "Our Answer", font_size=40,
             color=ACCENT, bold=True)

tf = add_text_box(slide, 1.0, 1.6, 5.5, 1, "Just dump it.",
                  font_size=32, color=WHITE, bold=True)
add_para(tf, "One input. That's all.", font_size=20, color=DIM,
         space_before=4)

# Code-like input box
add_shape_rect(slide, 1.0, 3.2, 5.5, 2.2, BG_MEDIUM)
tf = add_text_box(slide, 1.3, 3.4, 5, 1.8,
                  '\u201c Meeting with Sarah Friday 2pm,',
                  font_size=18, color=ACCENT2, font_name="Courier New")
add_para(tf, '  need to buy groceries,', font_size=18, color=ACCENT2,
         font_name="Courier New")
add_para(tf, '  oh and I finished the report \u201d', font_size=18,
         color=ACCENT2, font_name="Courier New")

# Right side: what AI does
add_shape_rect(slide, 7.2, 1.6, 5.3, 4.5, BG_MEDIUM)
tf = add_text_box(slide, 7.6, 1.8, 4.5, 0.5, "AI automatically:",
                  font_size=20, color=ACCENT, bold=True)
add_para(tf, "", font_size=10)
add_para(tf, "\u2713  Creates event: Sarah \u2014 Fri 2pm",
         font_size=19, color=ACCENT2, space_before=12)
add_para(tf, "\u2713  Creates task: Buy groceries",
         font_size=19, color=ACCENT2, space_before=12)
add_para(tf, "\u2713  Marks report as done!",
         font_size=19, color=ACCENT2, space_before=12)
add_para(tf, "", font_size=14)
add_para(tf, "No categorizing. No date pickers.",
         font_size=18, color=DIM, space_before=12)
add_para(tf, "No switching between apps.",
         font_size=18, color=DIM)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: OTTD
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.6, 10, 0.8,
             "OTTD \u2014 One Thing To Do", font_size=40,
             color=ACCENT, bold=True)

tf = add_text_box(slide, 1.0, 1.6, 7, 0.8,
                  "Instead of 20 tasks, we recommend just one.",
                  font_size=24, color=WHITE)

# Signal cards
signals = [
    ("\u23F0", "Time", "9am Monday\n\u2192 deep focus work", 1.0),
    ("\U0001F4CD", "Location", "Near grocery store\n\u2192 buy groceries now", 3.55),
    ("\U0001F4C5", "Calendar", "Meeting in 30min\n\u2192 pick something quick", 6.1),
    ("\u26a1", "Energy", "4 hours straight\n\u2192 take a break", 8.65),
    ("\U0001F60C", "Mood", "Feeling stressed\n\u2192 lighter task first", 11.2),
]

for emoji, title, desc, left in signals:
    add_shape_rect(slide, left, 2.8, 2.2, 3.5, BG_MEDIUM)
    add_text_box(slide, left, 2.95, 2.2, 0.6, emoji,
                 font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + 0.1, 3.7, 2.0, 0.5, title,
                 font_size=20, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + 0.1, 4.3, 2.0, 1.5, desc,
                 font_size=15, color=LIGHT, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.4)

tf = add_text_box(slide, 1.0, 6.6, 12, 0.5,
                  "You never feel overwhelmed. You always know what to do next.",
                  font_size=20, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: Wellness-First
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.6, 11, 0.8,
             "Not Just Productive. Healthy.", font_size=40,
             color=ACCENT, bold=True)

tf = add_text_box(slide, 1.0, 1.6, 11, 0.8,
                  "CAMP doesn't just optimize output. It protects you.",
                  font_size=22, color=LIGHT)

# Wellness recommendation examples
examples = [
    ("\U0001F33F", "You've been going hard all morning.\nTake a break. Maybe grab a juice \u2014\ngot any at home?"),
    ("\U0001F319", "It's 11pm. That task can wait.\nHow about winding down with something light?"),
    ("\u2615", "Great job finishing 5 tasks today!\nYou've done enough. Rest \u2014 tomorrow's light."),
]

for i, (icon, text) in enumerate(examples):
    left = 1.0 + i * 3.9
    add_shape_rect(slide, left, 3.0, 3.5, 3.0, BG_MEDIUM)
    add_text_box(slide, left, 3.15, 3.5, 0.6, icon,
                 font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + 0.3, 3.9, 2.9, 2.0, text,
                 font_size=17, color=ACCENT2, line_spacing=1.5)

tf = add_text_box(slide, 1.0, 6.4, 12, 0.6,
                  "The AI has intent: your wellness comes first.",
                  font_size=20, color=WARM, bold=True,
                  alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: Scenarios
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.6, 6, 0.8, "Scenarios", font_size=40,
             color=ACCENT, bold=True)

scenarios = [
    ("\u2600\ufe0f Morning Rush",
     "You dump:\n\"dentist at 3, pick up package,\nfinish slides, call mom\"",
     "CAMP creates 2 tasks + 1 event.\nRecommends: \"Finish slides first \u2014\nyou have a clear block until noon.\""),
    ("\U0001F319 After-Hours Grind",
     "It's 9pm.\nYou've completed 6 tasks today.",
     "CAMP says:\n\"Great day! You've done enough.\nRest \u2014 tomorrow's a light day.\""),
    ("\U0001F6B6 Walking Outside",
     "GPS detects you're near\ndowntown.",
     "CAMP nudges:\n\"You're near the post office \u2014\npick up that package now?\""),
]

for i, (title, situation, response) in enumerate(scenarios):
    top = 1.7 + i * 1.9
    add_shape_rect(slide, 1.0, top, 11.5, 1.65, BG_MEDIUM)
    add_text_box(slide, 1.3, top + 0.15, 2.5, 0.4, title,
                 font_size=20, color=ACCENT, bold=True)
    add_text_box(slide, 1.3, top + 0.6, 4.0, 1.0, situation,
                 font_size=16, color=LIGHT, line_spacing=1.4)
    # Arrow
    add_text_box(slide, 5.8, top + 0.6, 0.8, 1.0, "\u2192",
                 font_size=28, color=DIM, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 6.8, top + 0.6, 5.2, 1.0, response,
                 font_size=16, color=ACCENT2, line_spacing=1.4)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: Vision — Wearables
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.6, 10, 0.8,
             "Vision: Context-Aware AI", font_size=40,
             color=ACCENT, bold=True)

# Phase timeline
add_shape_rect(slide, 1.0, 1.8, 5.5, 2.0, BG_MEDIUM)
tf = add_text_box(slide, 1.3, 1.95, 5, 0.4, "Phase 1 \u2014 Now",
                  font_size=22, color=ACCENT2, bold=True)
add_para(tf, "", font_size=6)
add_para(tf, "\u2022  Time & calendar awareness", font_size=18, color=LIGHT)
add_para(tf, "\u2022  Location context (GPS)", font_size=18, color=LIGHT)
add_para(tf, "\u2022  Task history & patterns", font_size=18, color=LIGHT)

add_shape_rect(slide, 7.0, 1.8, 5.5, 2.0, BG_MEDIUM)
tf = add_text_box(slide, 7.3, 1.95, 5, 0.4, "Phase 2 \u2014 Wearables",
                  font_size=22, color=WARM, bold=True)
add_para(tf, "", font_size=6)
add_para(tf, "\u2764\ufe0f  Heart rate \u2192 detect stress, suggest breaks",
         font_size=18, color=LIGHT)
add_para(tf, "\U0001F4A4  Sleep quality \u2192 adjust task difficulty",
         font_size=18, color=LIGHT)
add_para(tf, "\U0001F3C3  Activity level \u2192 match energy to tasks",
         font_size=18, color=LIGHT)

# Big quote
add_shape_rect(slide, 1.0, 4.5, 11.5, 2.2, BG_MEDIUM)
tf = add_text_box(slide, 1.5, 4.8, 10.5, 1.8,
                  "CAMP becomes a personal chief of staff",
                  font_size=28, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)
add_para(tf, "who actually cares about you.",
         font_size=28, color=ACCENT, bold=True,
         alignment=PP_ALIGN.CENTER, space_before=8)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: Architecture
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.4, 10, 0.8, "How It Works", font_size=40,
             color=ACCENT, bold=True)

# Flow: User Input -> AI -> Actions -> DB -> Recommendation
boxes = [
    (1.0, 1.8, 2.8, 1.6, "User Input", "text / voice / image\njust dump anything", ACCENT),
    (4.5, 1.8, 3.2, 1.6, "Gemini 2.0 Flash", "Multimodal AI\nFunction Calling", ACCENT2),
    (8.5, 1.8, 3.8, 1.6, "Auto-Classification", "task / event / record\nexpense / habit log", WARM),
]

for left, top, w, h, title, desc, title_color in boxes:
    add_shape_rect(slide, left, top, w, h, BG_MEDIUM)
    add_text_box(slide, left + 0.2, top + 0.15, w - 0.4, 0.4, title,
                 font_size=18, color=title_color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + 0.2, top + 0.65, w - 0.4, 0.9, desc,
                 font_size=15, color=LIGHT, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.4)

# Arrows
add_text_box(slide, 3.85, 2.2, 0.6, 0.6, "\u2192", font_size=32,
             color=DIM, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 7.75, 2.2, 0.6, 0.6, "\u2192", font_size=32,
             color=DIM, alignment=PP_ALIGN.CENTER)

# Bottom row
add_shape_rect(slide, 1.0, 4.0, 5.5, 1.5, BG_MEDIUM)
tf = add_text_box(slide, 1.3, 4.1, 5, 0.4, "Supabase (PostgreSQL)",
                  font_size=18, color=ACCENT, bold=True)
add_para(tf, "Auth + RLS + Realtime + Storage", font_size=15, color=LIGHT)
add_para(tf, "Row-Level Security per user", font_size=15, color=LIGHT)

add_shape_rect(slide, 7.0, 4.0, 5.5, 1.5, BG_MEDIUM)
tf = add_text_box(slide, 7.3, 4.1, 5, 0.4, "Inngest (Background Jobs)",
                  font_size=18, color=ACCENT, bold=True)
add_para(tf, "\u2022 Auto-proceed: 2min silence \u2192 AI acts",
         font_size=15, color=LIGHT)
add_para(tf, "\u2022 Daily summary cron at midnight",
         font_size=15, color=LIGHT)

# OTTD box
add_shape_rect(slide, 3.5, 6.0, 6.3, 1.2, RGBColor(0x1E, 0x40, 0x6E))
tf = add_text_box(slide, 3.8, 6.1, 5.7, 0.4, "OTTD Recommendation Engine",
                  font_size=20, color=ACCENT, bold=True,
                  alignment=PP_ALIGN.CENTER)
add_para(tf, "Context scoring: time + location + energy + calendar \u2192 one task",
         font_size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: Tech Stack
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.6, 6, 0.8, "Tech Stack", font_size=40,
             color=ACCENT, bold=True)

# Left: Stack table
stack = [
    ("Frontend", "Next.js 16 (App Router) on Vercel"),
    ("AI Engine", "Gemini 2.0 Flash \u2014 multimodal, function calling"),
    ("Database", "Supabase PostgreSQL + Realtime + Auth"),
    ("Background", "Inngest (delayed jobs, cron)"),
    ("Recommendation", "Context scoring: time + location + energy"),
]

for i, (layer, tech) in enumerate(stack):
    top = 1.7 + i * 0.9
    add_shape_rect(slide, 1.0, top, 2.2, 0.7, BG_MEDIUM)
    add_text_box(slide, 1.1, top + 0.1, 2.0, 0.5, layer,
                 font_size=16, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 3.4, top + 0.1, 4.0, 0.5, tech,
                 font_size=16, color=LIGHT)

# Right: Key AI Behaviors
add_shape_rect(slide, 7.8, 1.7, 4.7, 4.5, BG_MEDIUM)
tf = add_text_box(slide, 8.1, 1.85, 4.2, 0.4, "Key AI Behaviors",
                  font_size=20, color=ACCENT, bold=True)

behaviors = [
    ("Classify", "One dump \u2192 multiple structured items"),
    ("Respond", "Conversational, 1-2 sentences, never overwhelming"),
    ("Nudge", "Auto-proceed after 2min silence (server-side)"),
    ("Heartbeat", "Daily summary, proactive habit suggestions"),
]

for label, desc in behaviors:
    add_para(tf, "", font_size=8)
    add_para(tf, f"\u25b6  {label}", font_size=17, color=ACCENT2,
             bold=True, space_before=8)
    add_para(tf, f"   {desc}", font_size=15, color=LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: Closing
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape_rect(slide, 0, 0, 0.15, 7.5, ACCENT)

add_text_box(slide, 1.5, 1.8, 10, 1.2, "CAMP", font_size=80,
             color=ACCENT, bold=True, font_name="Calibri",
             alignment=PP_ALIGN.CENTER)
tf = add_text_box(slide, 1.5, 3.5, 10, 0.8, "Calm Productivity.",
                  font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "", font_size=20)
add_para(tf, "Stop managing tasks.", font_size=28, color=DIM,
         alignment=PP_ALIGN.CENTER, space_before=12)
add_para(tf, "Start living.", font_size=28, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 5.8, 10, 0.6, "\u2014 Demo \u2014",
             font_size=24, color=ACCENT,
             alignment=PP_ALIGN.CENTER)


# Save
output_path = "/Users/edward/Projects/GitHub/writeflow/CAMP_Deck.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
